"""python-pptx based slide deck builder.

The Stage 1 PPT path (per docs/risk-reports/ppt_strategy_decision.md) is:
  outline → PptxBuilder.build → .pptx bytes → upload to Feishu Drive

This module owns step 1-2 only — turning a `SlideSchema` list into a
deterministic .pptx byte stream. It intentionally has zero network or
Feishu coupling so it can be unit-tested without any infra.
"""

import contextlib
import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.schemas.artifacts import SlideSchema
from app.schemas.enums import SlideLayout

_TITLE_FONT_PT = 32
_BULLET_FONT_PT = 18
_SUBTITLE_FONT_PT = 22

# python-pptx default 16:9 layout indexes (built-in template):
#   0 = Title (cover)
#   1 = Title and Content
#   5 = Title Only
#   6 = Blank
_LAYOUT_INDEX_BY_KIND: dict[SlideLayout, int] = {
    SlideLayout.cover: 0,
    SlideLayout.title_content: 1,
    SlideLayout.two_column: 1,  # python-pptx built-ins lack a clean 2-col; fall back
    SlideLayout.blank: 6,
    SlideLayout.section_header: 5,
}


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))  # type: ignore[no-untyped-call]


class PptxBuilder:
    """Pure builder: SlideSchema list → PowerPoint bytes."""

    def build(
        self,
        title: str,
        slides: list[SlideSchema],
        subtitle: str = "",
        token_name: str = "minimal",
    ) -> bytes:
        """Render a deck. Returns raw .pptx bytes (zip-format).

        token_name maps to a DesignToken preset (see app.services.design_tokens).
        Font sizes and text colour from the token are applied to all slides.
        """
        from app.services.design_tokens import get_preset

        try:
            token = get_preset(token_name)
        except KeyError:
            from app.services.design_tokens import get_preset as _gp

            token = _gp("minimal")

        self._token = token
        prs: Any = Presentation()

        # Cover slide is always first; if the outline already starts with a cover,
        # we just use that one and skip auto-injection.
        first_is_cover = bool(slides) and slides[0].layout == SlideLayout.cover
        if not first_is_cover:
            self._add_cover(prs, title, subtitle)

        for slide in slides:
            self._add_slide(prs, slide, deck_title=title)

        out = io.BytesIO()
        prs.save(out)
        return out.getvalue()

    def _add_cover(self, prs: Any, title: str, subtitle: str) -> None:
        layout = prs.slide_layouts[_LAYOUT_INDEX_BY_KIND[SlideLayout.cover]]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
            self._style_shape(slide.shapes.title, self._token.font_size_title)
        if subtitle and len(slide.placeholders) > 1:
            slide.placeholders[1].text = subtitle
            self._style_shape(slide.placeholders[1], _SUBTITLE_FONT_PT)

    def _add_slide(self, prs: Any, schema: SlideSchema, deck_title: str) -> None:
        layout_idx = _LAYOUT_INDEX_BY_KIND.get(schema.layout, 1)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)

        # Cover layout uses placeholder[1] for subtitle, others use it for body.
        if schema.layout == SlideLayout.cover:
            if slide.shapes.title is not None:
                slide.shapes.title.text = schema.title
                self._style_shape(slide.shapes.title, self._token.font_size_title)
            if schema.bullets and len(slide.placeholders) > 1:
                slide.placeholders[1].text = schema.bullets[0]
                self._style_shape(slide.placeholders[1], _SUBTITLE_FONT_PT)
            return

        if slide.shapes.title is not None:
            slide.shapes.title.text = schema.title
            self._style_shape(slide.shapes.title, self._token.font_size_title)

        body_ph = self._find_body_placeholder(slide)
        if body_ph is not None and schema.bullets:
            tf = body_ph.text_frame
            tf.clear()
            for i, bullet in enumerate(schema.bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.level = 0
                for run in p.runs:
                    run.font.size = Pt(self._token.font_size_body)
                    run.font.color.rgb = _hex_to_rgb(self._token.text_color)
        elif body_ph is None and schema.bullets:
            self._add_text_box(slide, schema.bullets, self._token.font_size_body)

        if schema.speaker_notes:
            slide.notes_slide.notes_text_frame.text = schema.speaker_notes

    def _style_shape(self, shape: Any, size_pt: int) -> None:
        """Apply token text colour and the given font size to all runs in a shape."""
        try:
            tf = shape.text_frame
        except AttributeError:
            return
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size_pt)
                with contextlib.suppress(Exception):
                    run.font.color.rgb = _hex_to_rgb(self._token.text_color)

    @staticmethod
    def _find_body_placeholder(slide: Any) -> Any:
        """Return the first content placeholder that isn't the title."""
        title_id = (
            slide.shapes.title.placeholder_format.idx if slide.shapes.title is not None else None
        )
        for ph in slide.placeholders:
            if title_id is not None and ph.placeholder_format.idx == title_id:
                continue
            return ph
        return None

    def _add_text_box(self, slide: Any, bullets: list[str], font_size_pt: int) -> None:
        """Fallback: add a textbox when no body placeholder exists (e.g. blank layout)."""
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            for run in p.runs:
                run.font.size = Pt(font_size_pt)
                with contextlib.suppress(Exception):
                    run.font.color.rgb = _hex_to_rgb(self._token.text_color)

    @staticmethod
    def _set_run_size(shape: Any, size_pt: int) -> None:
        try:
            tf = shape.text_frame
        except AttributeError:
            return
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(size_pt)
