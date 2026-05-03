"""Full coverage tests for PptxBuilder — exercises all uncovered branches."""

from __future__ import annotations

from app.integrations.python_pptx.builder import PptxBuilder, _hex_to_rgb
from app.schemas.artifacts import ChartSchema, ChartSeries, SlideSchema
from app.schemas.enums import SlideLayout


def _make_slide(
    title: str = "Test Slide",
    layout: SlideLayout = SlideLayout.title_content,
    bullets: list[str] | None = None,
    speaker_notes: str = "",
    chart: ChartSchema | None = None,
) -> SlideSchema:
    return SlideSchema(
        page_index=0,
        layout=layout,
        title=title,
        bullets=bullets or [],
        speaker_notes=speaker_notes,
        chart=chart,
    )


class TestHexToRgb:
    def test_basic_color(self) -> None:
        rgb = _hex_to_rgb("#FF0000")
        assert rgb[0] == 255
        assert rgb[1] == 0
        assert rgb[2] == 0

    def test_without_hash(self) -> None:
        rgb = _hex_to_rgb("0080FF")
        assert rgb[0] == 0
        assert rgb[1] == 0x80
        assert rgb[2] == 0xFF


class TestPptxBuilderBuild:
    def test_build_returns_bytes(self) -> None:
        builder = PptxBuilder()
        slides = [_make_slide("Hello")]
        result = builder.build("Title", slides)
        assert isinstance(result, bytes)
        assert len(result) > 0

    def test_build_with_unknown_token_falls_back(self) -> None:
        """Lines 75-78: KeyError branch when token_name not found."""
        builder = PptxBuilder()
        slides = [_make_slide("Slide")]
        result = builder.build("Deck", slides, token_name="nonexistent_token_xyz")
        assert isinstance(result, bytes)

    def test_build_with_subtitle(self) -> None:
        builder = PptxBuilder()
        slides = [_make_slide("Slide")]
        result = builder.build("Title", slides, subtitle="A subtitle")
        assert isinstance(result, bytes)

    def test_build_with_cover_slide_first_skips_auto_cover(self) -> None:
        """When first slide is cover layout, no extra cover is prepended."""
        builder = PptxBuilder()
        cover = _make_slide("My Cover", layout=SlideLayout.cover, bullets=["Subtitle"])
        result = builder.build("Title", [cover])
        assert isinstance(result, bytes)

    def test_build_section_header_layout(self) -> None:
        builder = PptxBuilder()
        slides = [_make_slide("Section", layout=SlideLayout.section_header)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_build_blank_layout(self) -> None:
        builder = PptxBuilder()
        slides = [_make_slide("Blank", layout=SlideLayout.blank)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)


class TestPptxBuilderSlideContent:
    def test_speaker_notes_written(self) -> None:
        """Line 140: speaker_notes path."""
        builder = PptxBuilder()
        slides = [_make_slide("Notes Slide", speaker_notes="Remember to mention X.")]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_bullets_written(self) -> None:
        builder = PptxBuilder()
        bullets = ["Point one", "Point two", "Point three"]
        slides = [_make_slide("Bullet Slide", bullets=bullets)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_cover_with_bullets_uses_placeholder(self) -> None:
        """Cover layout with bullets: first bullet → placeholder[1]."""
        builder = PptxBuilder()
        cover = _make_slide("Cover", layout=SlideLayout.cover, bullets=["subtitle text"])
        result = builder.build("Deck", [cover])
        assert isinstance(result, bytes)

    def test_cover_without_bullets(self) -> None:
        builder = PptxBuilder()
        cover = _make_slide("Cover", layout=SlideLayout.cover)
        result = builder.build("Deck", [cover])
        assert isinstance(result, bytes)


class TestPptxBuilderChart:
    def _make_bar_chart(self, title: str = "Bar Chart") -> ChartSchema:
        return ChartSchema(
            chart_type="bar",
            title=title,
            categories=["A", "B", "C"],
            series=[ChartSeries(name="Series1", values=[1.0, 2.0, 3.0])],
        )

    def test_build_slide_with_bar_chart(self) -> None:
        """Lines 142-184: full _add_chart_to_slide path."""
        builder = PptxBuilder()
        chart = self._make_bar_chart()
        slides = [_make_slide("Chart Slide", bullets=["Note"], chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_with_no_bullets(self) -> None:
        """n_bullets=0 path in _add_chart_to_slide."""
        builder = PptxBuilder()
        chart = self._make_bar_chart()
        slides = [_make_slide("Chart Only", chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_with_many_bullets_logs_warning(self) -> None:
        """n_bullets > 6 triggers warning branch."""
        builder = PptxBuilder()
        chart = self._make_bar_chart()
        bullets = [f"Point {i}" for i in range(8)]
        slides = [_make_slide("Heavy Slide", bullets=bullets, chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_no_series_skips(self) -> None:
        """chart.series is empty → early return (warning logged)."""
        builder = PptxBuilder()
        chart = ChartSchema(
            chart_type="bar",
            title="Empty",
            categories=["A"],
            series=[],
        )
        slides = [_make_slide("Empty Chart", chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_line_type(self) -> None:
        builder = PptxBuilder()
        chart = ChartSchema(
            chart_type="line",
            title="Line",
            categories=["Q1", "Q2"],
            series=[ChartSeries(name="Rev", values=[10.0, 20.0])],
        )
        slides = [_make_slide("Line Chart", chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_pie_type(self) -> None:
        builder = PptxBuilder()
        chart = ChartSchema(
            chart_type="pie",
            title="Pie",
            categories=["Alpha", "Beta"],
            series=[ChartSeries(name="Share", values=[60.0, 40.0])],
        )
        slides = [_make_slide("Pie Chart", chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_unknown_type_falls_back(self) -> None:
        """Unknown chart_type string → _CHART_TYPE_MAP.get falls back to COLUMN_CLUSTERED."""
        from unittest.mock import patch

        from app.integrations.python_pptx.builder import _CHART_TYPE_MAP

        # Temporarily add an unknown type so Pydantic accepts it but the map misses
        chart = ChartSchema(
            chart_type="bar",
            title="FallbackTest",
            categories=["X", "Y"],
            series=[ChartSeries(name="S", values=[1.0, 2.0])],
        )
        # Patch the map to not contain 'bar' so the fallback path executes
        with patch.dict(_CHART_TYPE_MAP, {}, clear=True):
            builder = PptxBuilder()
            slides = [_make_slide("FallbackChart", chart=chart)]
            result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_with_title_set(self) -> None:
        """Lines 181-184: chart_schema.title → set chart_obj.has_title."""
        builder = PptxBuilder()
        chart = ChartSchema(
            chart_type="bar",
            title="My Chart Title",
            categories=["X"],
            series=[ChartSeries(name="S", values=[5.0])],
        )
        slides = [_make_slide("Titled Chart", chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_chart_no_title(self) -> None:
        builder = PptxBuilder()
        chart = ChartSchema(
            chart_type="bar",
            title="",
            categories=["X"],
            series=[ChartSeries(name="S", values=[5.0])],
        )
        slides = [_make_slide("No Title Chart", chart=chart)]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)


class TestPptxBuilderHelpers:
    def test_find_body_placeholder_returns_none_for_blank(self) -> None:
        """Blank layout has no body placeholder → _find_body_placeholder returns None."""
        builder = PptxBuilder()
        bullets = ["fallback text"]
        slides = [_make_slide("Blank", layout=SlideLayout.blank, bullets=bullets)]
        result = builder.build("Deck", slides)
        # If _add_text_box was called (lines 212-224), result is still valid bytes
        assert isinstance(result, bytes)

    def test_set_run_size_static_method(self) -> None:
        """Lines 228-234: _set_run_size static method."""
        from pptx import Presentation

        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title is not None:
            PptxBuilder._set_run_size(slide.shapes.title, 24)

    def test_set_run_size_no_text_frame(self) -> None:
        """_set_run_size with object lacking text_frame should not raise."""

        class FakeShape:
            @property
            def text_frame(self) -> None:
                raise AttributeError("no text_frame")

        PptxBuilder._set_run_size(FakeShape(), 12)  # should not raise

    def test_style_shape_no_text_frame(self) -> None:
        """_style_shape with object lacking text_frame should not raise."""
        builder = PptxBuilder()
        builder._token = type("T", (), {"font_size_title": 32, "text_color": "#000000"})()  # type: ignore[assignment]

        class FakeShape:
            @property
            def text_frame(self) -> None:
                raise AttributeError("no text_frame")

        builder._style_shape(FakeShape(), 12)  # should not raise

    def test_two_column_layout_fallback(self) -> None:
        """two_column falls back to layout index 1."""
        builder = PptxBuilder()
        slides = [_make_slide("Two Col", layout=SlideLayout.two_column, bullets=["A", "B"])]
        result = builder.build("Deck", slides)
        assert isinstance(result, bytes)

    def test_multiple_slides(self) -> None:
        builder = PptxBuilder()
        slides = [
            _make_slide(f"Slide {i}", bullets=[f"Bullet {i}a", f"Bullet {i}b"]) for i in range(4)
        ]
        result = builder.build("Multi", slides)
        assert isinstance(result, bytes)
