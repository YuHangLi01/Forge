"""Unit tests for PPTService + PptxBuilder + outline loader.

PPTService is async glue around a synchronous builder, so most tests just
verify that the bytes coming out of the builder are a valid .pptx (zip)
and that the slide count / title / bullets land where expected.

Reading the resulting .pptx back via python-pptx is the cheapest content
oracle — no need to parse XML by hand.
"""

import io
import json
from pathlib import Path

import pytest
from pptx import Presentation as PptxOpen

from app.schemas.artifacts import PPTArtifact, SlideSchema
from app.schemas.enums import SlideLayout
from app.services.ppt_outline_loader import load_outline
from app.services.ppt_service import PPTService

_FIXTURE_DIR = Path(__file__).resolve().parent / "fixtures" / "outlines"


def _make_slides() -> list[SlideSchema]:
    return [
        SlideSchema(
            page_index=0,
            layout=SlideLayout.cover,
            title="Demo Deck",
            bullets=["sub-title text"],
        ),
        SlideSchema(
            page_index=1,
            layout=SlideLayout.title_content,
            title="第一页内容",
            bullets=["要点 A", "要点 B", "要点 C"],
        ),
        SlideSchema(
            page_index=2,
            layout=SlideLayout.title_content,
            title="第二页内容",
            bullets=["only one bullet"],
        ),
    ]


def _open_deck(pptx_bytes: bytes):  # type: ignore[no-untyped-def]
    return PptxOpen(io.BytesIO(pptx_bytes))


@pytest.mark.asyncio
async def test_build_pptx_bytes_returns_zip_bytes() -> None:
    bytes_ = await PPTService().build_pptx_bytes("Demo", _make_slides())
    assert bytes_[:2] == b"PK", "pptx is a zip; should start with 'PK'"
    assert len(bytes_) > 5_000


@pytest.mark.asyncio
async def test_build_correct_slide_count_when_outline_starts_with_cover() -> None:
    bytes_ = await PPTService().build_pptx_bytes("Demo", _make_slides())
    deck = _open_deck(bytes_)
    # outline already has cover at index 0 → builder should NOT inject another
    assert len(deck.slides) == 3


@pytest.mark.asyncio
async def test_build_auto_injects_cover_when_outline_lacks_one() -> None:
    bullet_slides = [s for s in _make_slides() if s.layout != SlideLayout.cover]
    bytes_ = await PPTService().build_pptx_bytes("Demo", bullet_slides, subtitle="auto-cover")
    deck = _open_deck(bytes_)
    # 1 auto cover + 2 content = 3
    assert len(deck.slides) == 3


@pytest.mark.asyncio
async def test_build_renders_title_and_bullets() -> None:
    bytes_ = await PPTService().build_pptx_bytes("Demo", _make_slides())
    deck = _open_deck(bytes_)
    titles = [s.shapes.title.text for s in deck.slides if s.shapes.title is not None]
    assert "Demo Deck" in titles
    assert "第一页内容" in titles

    # Find the slide with title "第一页内容" and check bullets
    first_content = next(
        s for s in deck.slides if s.shapes.title and s.shapes.title.text == "第一页内容"
    )
    body_text = "\n".join(
        p.text
        for ph in first_content.placeholders
        if ph.placeholder_format.idx != first_content.shapes.title.placeholder_format.idx
        for p in ph.text_frame.paragraphs
    )
    assert "要点 A" in body_text
    assert "要点 B" in body_text
    assert "要点 C" in body_text


@pytest.mark.asyncio
async def test_build_handles_empty_bullets() -> None:
    slides = [
        SlideSchema(
            page_index=0,
            layout=SlideLayout.title_content,
            title="只有标题，没要点",
            bullets=[],
        )
    ]
    bytes_ = await PPTService().build_pptx_bytes("Demo", slides)
    # 1 auto cover + 1 content
    deck = _open_deck(bytes_)
    assert len(deck.slides) == 2


@pytest.mark.asyncio
async def test_create_from_outline_returns_artifact_without_adapter() -> None:
    artifact = await PPTService().create_from_outline("Demo", _make_slides())
    assert isinstance(artifact, PPTArtifact)
    assert artifact.title == "Demo"
    assert artifact.ppt_id == ""  # no adapter → no upload → empty
    assert artifact.share_url == ""
    assert len(artifact.slides) == 3


@pytest.mark.asyncio
async def test_patch_slide_raises_not_implemented() -> None:
    svc = PPTService()
    with pytest.raises(NotImplementedError, match="Stage 2"):
        await svc.patch_slide("doc-token", 0, _make_slides()[0])


# ---------- outline loader ----------


@pytest.mark.parametrize("name", [p.stem for p in _FIXTURE_DIR.glob("*.json")])
def test_outline_fixtures_load(name: str) -> None:
    title, subtitle, slides = load_outline(_FIXTURE_DIR / f"{name}.json")
    assert title, f"{name}: title empty"
    assert len(slides) >= 5, f"{name}: too few slides ({len(slides)})"
    # First slide should be a cover (type='title' in fixture)
    assert slides[0].layout == SlideLayout.cover, f"{name}: first slide not cover"


def test_outline_loader_rejects_malformed() -> None:
    bad = '{"title": "x", "slides": "not-a-list"}'
    p = _FIXTURE_DIR.parent / "_bad.json"
    p.write_text(bad)
    try:
        with pytest.raises(ValueError, match="must be a list"):
            load_outline(p)
    finally:
        p.unlink()


@pytest.mark.asyncio
async def test_build_a_real_fixture_end_to_end() -> None:
    """Round-trip: load outline JSON → build pptx → verify slide count."""
    fixture = _FIXTURE_DIR / "01_requirements.json"
    title, subtitle, slides = load_outline(fixture)
    bytes_ = await PPTService().build_pptx_bytes(title, slides, subtitle=subtitle)
    deck = _open_deck(bytes_)

    # outline has 11 slides (1 title + 10 content); builder respects them
    raw = json.loads(fixture.read_text(encoding="utf-8"))
    expected_n = len(raw["slides"])
    assert len(deck.slides) == expected_n
